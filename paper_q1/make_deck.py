#!/usr/bin/env python
"""Build the presentation deck (PPTX) for the coded-audio XAI study.

Twenty 16:9 slides: problem statement, literature and gap, dataset and design,
methodology, results, conclusion. Every number on a slide is taken from
`paper_q1/main.tex` and `xai_ser/reports/FINDINGS.md`, which in turn read the
committed result files; nothing here is estimated.

This writes a NEW deck. The pre-existing HTML deck under
`xai_ser/outputs/presentation/` is left untouched.

    python3 paper_q1/make_deck.py
"""

from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from PIL import Image

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.font_manager import FontProperties

HERE = Path(__file__).resolve().parent
FIGS = HERE / "figs"
DEST = HERE / "WavVsMp3_deck.pptx"

# ----------------------------------------------------------------- palette
INK     = RGBColor(0x0D, 0x14, 0x18)
BODY    = RGBColor(0x24, 0x32, 0x38)
MUTED   = RGBColor(0x5A, 0x6B, 0x74)
FAINT   = RGBColor(0x8C, 0x9B, 0xA2)
RULE    = RGBColor(0xD4, 0xDC, 0xE0)
HAIR    = RGBColor(0xE7, 0xEC, 0xEE)
PAPER   = RGBColor(0xFF, 0xFF, 0xFF)
TINT    = RGBColor(0xEF, 0xF4, 0xF7)
WAV     = RGBColor(0x2E, 0x6F, 0x9E)   # methodology / uncompressed
MP3     = RGBColor(0xC4, 0x46, 0x2E)   # results / codec damage
VOICE   = RGBColor(0x3F, 0x8F, 0x5B)   # conclusion / recovery
AMBER   = RGBColor(0xB0, 0x7A, 0x1E)

DISPLAY = "Helvetica Neue"
MONO    = "Menlo"

# ----------------------------------------------------------------- geometry
SW, SH  = 13.333, 7.5
ML, MR  = 0.66, 0.66
CW      = SW - ML - MR            # 12.013 in of live width
TITLE_Y = 0.46
RULE_Y  = 1.30
BODY_Y  = 1.56
BODY_B  = 6.80                    # bottom of the body region
FOOT_Y  = 6.98

SECTION_COLOR = {
    "intro": MUTED, "method": WAV, "result": MP3, "close": VOICE,
}

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

_slides: list = []


# ------------------------------------------------------------ oxml helpers
def _letter_space(run, points: float) -> None:
    """Track-out a run. PowerPoint stores this as 1/100 pt on the run props."""
    run.font._rPr.set("spc", str(int(points * 100)))


_LN_SUCC = {
    "lnL": ("a:lnR", "a:lnT", "a:lnB", "a:lnTlToBr", "a:lnBlToTr", "a:cell3D",
            "a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill",
            "a:grpFill", "a:headers", "a:extLst"),
    "lnR": ("a:lnT", "a:lnB", "a:lnTlToBr", "a:lnBlToTr", "a:cell3D", "a:noFill",
            "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill",
            "a:headers", "a:extLst"),
    "lnT": ("a:lnB", "a:lnTlToBr", "a:lnBlToTr", "a:cell3D", "a:noFill",
            "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill",
            "a:headers", "a:extLst"),
    "lnB": ("a:lnTlToBr", "a:lnBlToTr", "a:cell3D", "a:noFill", "a:solidFill",
            "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill", "a:headers",
            "a:extLst"),
}


def cell_rule(cell, edge: str, color: RGBColor, pt: float) -> None:
    """Draw one edge of a table cell. Only horizontal rules are used, so the
    tables read like the booktabs tables in the paper."""
    from pptx.oxml.xmlchemy import OxmlElement

    tcPr = cell._tc.get_or_add_tcPr()
    for old in tcPr.findall(qn(f"a:{edge}")):
        tcPr.remove(old)
    ln = OxmlElement(f"a:{edge}")
    ln.set("w", str(Pt(pt)))
    ln.set("cap", "flat")
    fill = OxmlElement("a:solidFill")
    clr = OxmlElement("a:srgbClr")
    clr.set("val", f"{color}")
    fill.append(clr)
    ln.append(fill)
    tcPr.insert_element_before(ln, *_LN_SUCC[edge])


# ------------------------------------------------------------ text helpers
_MARKUP = re.compile(r"(\*\*.+?\*\*|\*[^*`]+?\*|`.+?`)")

_MFIG = plt.figure(figsize=(1, 1), dpi=100)
_MREN = _MFIG.canvas.get_renderer()
WARN: list = []


def _tw(t, name, size, bold):
    prop = FontProperties(family=name, size=size,
                          weight="bold" if bold else "normal")
    return _MREN.get_text_width_height_descent(t, prop, False)[0] / 100.0


def _split(text, size, bold=False):
    out = []
    for chunk in _MARKUP.split(text):
        if not chunk:
            continue
        if chunk.startswith("**"):
            out.append((chunk[2:-2], DISPLAY, size, True))
        elif chunk.startswith("*"):
            out.append((chunk[1:-1], DISPLAY, size, bold))
        elif chunk.startswith("`"):
            out.append((chunk[1:-1], MONO, size * 0.90, False))
        else:
            out.append((chunk, DISPLAY, size, bold))
    return out


def nlines(text, width, size, bold=False):
    """Line count after greedy wrapping, measured with the real font metrics."""
    lines, x = 1, 0.0
    for t, f, sz, b in _split(text, size, bold):
        for k, wd in enumerate(t.split(" ")):
            if not wd and k:
                continue
            tok = wd if x == 0 else " " + wd
            w = _tw(tok, f, sz, b)
            if x + w > width and x > 0:
                lines += 1
                x = _tw(wd, f, sz, b)
            else:
                x += w
    return lines


def check(tag, bottom, limit=None):
    limit = BODY_B if limit is None else limit
    if bottom > limit + 0.005:
        WARN.append(f"{tag}: overflows by {bottom - limit:.2f} in "
                    f"(ends {bottom:.2f}, limit {limit:.2f})")
    return bottom


def write(para, text: str, size: float, color: RGBColor = BODY,
          bold: bool = False, italic: bool = False, font: str = DISPLAY) -> None:
    """Render a mini-markup string into a paragraph.

    ``**bold**`` emphasises, and ``` `code` ``` switches to the monospace face
    used for descriptor names, so `spectral_contrast_6_mean` reads on a slide
    the way it reads in the paper.
    """
    for chunk in _MARKUP.split(text):
        if not chunk:
            continue
        r = para.add_run()
        if chunk.startswith("**"):
            r.text, r.font.bold, r.font.name = chunk[2:-2], True, font
            r.font.color.rgb = INK if color is BODY else color
        elif chunk.startswith("*"):
            r.text, r.font.name = chunk[1:-1], font
            r.font.bold, r.font.italic = bold, True
            r.font.color.rgb = color
            r.font.size = Pt(size)
            continue
        elif chunk.startswith("`"):
            r.text, r.font.name = chunk[1:-1], MONO
            r.font.size = Pt(size * 0.90)
            r.font.color.rgb = color
            r.font.bold, r.font.italic = bold, italic
            continue
        else:
            r.text, r.font.bold, r.font.name = chunk, bold, font
            r.font.color.rgb = color
        r.font.size = Pt(size)
        r.font.italic = italic


def textbox(slide, x, y, w, h, *, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def line(slide, x, y, w, color=RULE, pt=0.75):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Pt(pt))
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def panel(slide, x, y, w, h, fill=TINT, edge=None, radius=0.035):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if edge is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = edge
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


# ------------------------------------------------------------ slide chrome
def slide(title: str, eyebrow: str, section: str = "result"):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER
    accent = SECTION_COLOR[section]

    tf = textbox(s, ML, TITLE_Y, CW, 0.30)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = eyebrow.upper()
    r.font.size, r.font.bold, r.font.name = Pt(10), True, DISPLAY
    r.font.color.rgb = accent
    _letter_space(r, 1.4)

    tf = textbox(s, ML, TITLE_Y + 0.30, CW, 0.62)
    p = tf.paragraphs[0]
    write(p, title, 24, INK, bold=True)

    line(s, ML, RULE_Y, CW, HAIR, 0.75)
    line(s, ML, RULE_Y, 1.10, accent, 1.6)

    _slides.append(s)
    return s


def footer_all() -> None:
    """Numbering is stamped once at the end so the totals are right."""
    n = len(_slides)
    for i, s in enumerate(_slides[1:], start=2):
        line(s, ML, FOOT_Y - 0.16, CW, HAIR, 0.5)
        tf = textbox(s, ML, FOOT_Y, CW * 0.72, 0.24)
        p = tf.paragraphs[0]
        write(p, "Null-Calibrated Feature Attribution in Explainable SER  ·  CREMA-D",
              9, FAINT)
        tf = textbox(s, ML + CW * 0.72, FOOT_Y, CW * 0.28, 0.24)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        write(p, f"{i} / {n}", 9, FAINT)


# ------------------------------------------------------------ content parts
def bullets(slide, items, x, y, w, *, size=14.0, gap=9.0, color=BODY,
            marker=True, accent=WAV, lead=1.16, tag=None, limit=None):
    """`items` are (level, text) pairs. Level 0 takes a square marker."""
    tf = textbox(slide, x, y, w, BODY_B - y)
    first = True
    for level, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(0 if p is tf.paragraphs[0] else gap)
        p.space_after = Pt(0)
        p.line_spacing = lead
        if level == 0:
            if marker:
                r = p.add_run()
                r.text = "\u2022  "
                r.font.size, r.font.name = Pt(size * 0.72), DISPLAY
                r.font.color.rgb = accent
            write(p, text, size, color)
        else:
            p.level = 1
            r = p.add_run()
            r.text = "     "
            r.font.size = Pt(size * 0.86)
            write(p, text, size * 0.86, MUTED)

    bottom = y
    for i, (level, text) in enumerate(items):
        sz = size if level == 0 else size * 0.86
        pre = "\u2022  " if (level == 0 and marker) else "     "
        n = nlines(pre + text, w, sz)
        bottom += (0 if i == 0 else gap / 72.0) + n * sz * 1.21 * lead / 72.0
    return check(tag or "bullets", bottom, limit)


def band_h(text, w, size=13.5):
    """Height a takeaway band needs for `text`, so the block above it can be
    given a matching limit."""
    n = nlines(text, w - 0.55, size)
    return max(0.62, n * size * 1.14 * 1.21 / 72.0 + 0.16)


def takeaway(slide, text, *, h=None, accent=MP3, size=13.5, x=ML, w=CW, tag=None):
    h = band_h(text, w, size) if h is None else h
    y = BODY_B - h
    panel(slide, x, y, w, h, TINT)
    line(slide, x, y, 0.055, accent, h * 72)
    tf = textbox(slide, x + 0.30, y, w - 0.55, h, anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    p.line_spacing = 1.14
    write(p, text, size, INK)
    return tf


def table(slide, rows, x, y, w, *, col_w=None, size=11.0, head_size=11.0,
          row_h=0.255, head_h=0.30, align=None, accent=INK, rule_rows=()):
    """A booktabs-style table: top rule, header rule, bottom rule, no verticals.

    `rows[0]` is the header. `rule_rows` names 0-based body rows that get a
    rule above them, which is how the paper separates its blocks.
    """
    nr, nc = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w),
                                Inches(head_h + row_h * (nr - 1)))
    tbl = gf.table
    tbl.first_row = False
    tbl.horz_banding = False

    if col_w:
        total = sum(col_w)
        for i, frac in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(w) * frac / total))
    tbl.rows[0].height = Inches(head_h)
    for i in range(1, nr):
        tbl.rows[i].height = Inches(row_h)

    align = align or (["l"] + ["c"] * (nc - 1))
    amap = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.background()
            cell.margin_left = cell.margin_right = Inches(0.055)
            cell.margin_top = cell.margin_bottom = Inches(0.018)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = amap[align[ci]]
            p.line_spacing = 1.0
            head = ri == 0
            write(p, str(val), head_size if head else size,
                  accent if head else BODY, bold=head)
            if head:
                cell_rule(cell, "lnT", INK, 1.1)
                cell_rule(cell, "lnB", INK, 0.75)
            elif ri == nr - 1:
                cell_rule(cell, "lnB", INK, 1.1)
            if ri in rule_rows:
                cell_rule(cell, "lnT", RULE, 0.75)
    return tbl


def picture(slide, name, x, y, w, h, *, caption=None):
    """Fit a figure inside a box, preserving aspect, centred."""
    path = FIGS / f"{name}.png"
    iw, ih = Image.open(path).size
    ar = iw / ih
    if w / h > ar:
        ph, pw = h, h * ar
    else:
        pw, ph = w, w / ar
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    slide.shapes.add_picture(str(path), Inches(px), Inches(py),
                             Inches(pw), Inches(ph))
    if caption:
        tf = textbox(slide, x, py + ph + 0.07, w, 0.32)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.1
        write(p, caption, 9.5, MUTED, italic=True)
    return px, py, pw, ph


def stat(slide, x, y, w, value, label, *, color=MP3, vsize=30, lsize=10.5):
    tf = textbox(slide, x, y, w, 0.52)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    write(p, value, vsize, color, bold=True)
    tf = textbox(slide, x, y + vsize * 1.30 / 72.0, w, 0.62)
    for j, ln in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = 1.12
        write(p, ln, lsize, MUTED)


def prose(slide, x, y, w, paras, *, size=12.0, color=BODY, gap=8.0,
          lead=1.18, tag=None, limit=None):
    """A measured multi-paragraph text block. `paras` are (text, size, color)
    or plain strings taking the defaults."""
    tf = textbox(slide, x, y, w, BODY_B - y)
    bottom = y
    for i, item in enumerate(paras):
        text, sz, col = item if isinstance(item, tuple) else (item, size, color)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0 if i == 0 else gap)
        p.line_spacing = lead
        write(p, text, sz, col)
        bottom += (0 if i == 0 else gap / 72.0) + \
            nlines(text, w, sz) * sz * 1.21 * lead / 72.0
    return check(tag or "prose", bottom, limit)


def eyebrow(slide, x, y, w, text, color=MUTED, size=10.0, track=1.2):
    tf = textbox(slide, x, y, w, 0.26)
    p = tf.paragraphs[0]
    write(p, text.upper(), size, color, bold=True)
    _letter_space(p.runs[-1], track)
    return y + 0.30


# =========================================================================
# 1 — title
# =========================================================================
s = prs.slides.add_slide(BLANK)
s.background.fill.solid()
s.background.fill.fore_color.rgb = PAPER
_slides.append(s)

panel(s, 0, 0, SW, 0.30, INK, radius=0.0)
line(s, ML, 1.62, 2.20, MP3, 2.2)

tf = textbox(s, ML, 1.86, CW, 2.30)
for i, ln in enumerate(["What the Model Learns", "from Coded Audio"]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.line_spacing = 1.06
    write(p, ln, 44, INK, bold=True)

tf = textbox(s, ML, 3.58, CW * 0.80, 0.90)
p = tf.paragraphs[0]
p.line_spacing = 1.24
write(p, "Null-Calibrated Feature Attribution in "
         "Explainable Speech Emotion Recognition", 19, MUTED)

line(s, ML, 4.50, CW, HAIR, 0.75)

tf = textbox(s, ML, 4.70, CW * 0.5, 0.60)
p = tf.paragraphs[0]
write(p, "Akshath R", 15, INK, bold=True)
p = tf.add_paragraph()
p.space_before = Pt(3)
write(p, "Independent researcher  ·  akshath.r333@gmail.com", 12, MUTED)

FACTS = [("7,442", "CREMA-D clips\n91 actors"),
         ("4", "coding conditions\n29,768 files"),
         ("436", "named acoustic\ndescriptors"),
         ("15", "classifiers\n+ 10 XAI methods")]
bw = CW / 4
for i, (v, l) in enumerate(FACTS):
    x = ML + i * bw
    if i:
        line(s, x, 5.72, 0.008, HAIR, 62)
    tf = textbox(s, x + 0.02, 5.66, bw - 0.20, 0.44)
    write(tf.paragraphs[0], v, 27, WAV if i % 2 == 0 else MP3, bold=True)
    tf = textbox(s, x + 0.02, 6.08, bw - 0.20, 0.62)
    for j, ln in enumerate(l.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.line_spacing = 1.16
        write(p, ln, 11, MUTED)

panel(s, 0, SH - 0.16, SW, 0.16, INK, radius=0.0)


# =========================================================================
# 2 — problem statement
# =========================================================================
s = slide("Models are validated on WAV and deployed on coded audio",
          "Introduction  ·  Problem statement", "intro")

bullets(s, [
    (0, "Benchmarks ship as **uncompressed WAV** and are read straight from "
        "disk. Almost nothing reaches an inference endpoint that way."),
    (1, "Contact-centre archives hold MP3. Streaming clients negotiate AAC or "
        "Opus. Broadcast pipelines resample and requantise."),
    (0, "The usual defence: every ingestion path decodes to linear PCM before "
        "features are extracted, so the delivery format looks like an "
        "implementation detail."),
    (0, "Where that defence has been tested, it has been tested on "
        "**accuracy** — and accuracy is now only half of what a deployed "
        "affective system is required to report."),
    (1, "A clinician, an auditor or a QA analyst is shown an attribution over "
        "acoustic evidence. That second output has its own robustness "
        "properties, and they do not follow from the first."),
], ML, BODY_Y, CW * 0.575, size=14.5, gap=11, accent=MUTED,
   tag="s2 left", limit=5.98)

px = ML + CW * 0.615
panel(s, px, BODY_Y - 0.04, CW * 0.385, 3.34, TINT)
prose(s, px + 0.30, BODY_Y + 0.20, CW * 0.385 - 0.60, [
    ("THE QUESTION", 10, MP3),
    ("Does perceptual coding change the **evidence** a speech-emotion model "
     "uses — not just how often it is right?", 17, INK),
    ("Answered with models whose evidence is a **named acoustic descriptor** "
     "rather than a latent dimension, so every attribution has a phonetic "
     "reading by construction.", 13, MUTED),
], gap=12, tag="s2 panel", limit=BODY_Y + 3.22)

takeaway(s, "Attribution fragility has been demonstrated under **adversarial** "
            "perturbation. A transcode is not adversarial. It is routine, it is "
            "everywhere, and it is almost never logged.",
         accent=MUTED, tag="s2 band")


# =========================================================================
# 3 — literature
# =========================================================================
s = slide("What the literature establishes, and where it stops",
          "Literature", "intro")

COLS = [
    ("Coding × speech tasks", WAV, [
        "**Reddy & Vijayarajan (2020)** — codec effects on SER are real, "
        "modest, and depend on both codec family and front end. Coding "
        "interacts with the representation rather than degrading uniformly.",
        "**Wu et al. (2024), Codec-SUPERB** — signal-domain fidelity metrics "
        "*fail to predict* whether task-relevant information survives. A codec "
        "can score well on reconstruction and still discard paralinguistics.",
    ]),
    ("Attribution fragility", MP3, [
        "**Ghorbani et al. (2019)** — perturbations below human perception, "
        "with identical predictions, still produce very different attributions.",
        "**Adebayo et al. (2018)** — some saliency methods are statistically "
        "independent of the model they claim to explain.",
        "**Krishna et al. (2022)** — methods disagree on the same fitted model, "
        "and practitioners resolve it on no principled basis.",
    ]),
    ("Interpretability for audio", VOICE, [
        "**Haunschmid et al. (2020)** — audio explanations must be "
        "*listenable*; image-segmentation logic on spectrograms yields "
        "components with no perceptual correlate. **Nasr et al. (2025)** extend "
        "the critique to speech emotion specifically.",
        "**Rudin (2019)** — in high-stakes settings, replace the post-hoc layer "
        "with an intrinsically interpretable model. Our leaderboard speaks to "
        "this directly, and does not support it.",
    ]),
]
cw = (CW - 0.62) / 3
for i, (head, col, points) in enumerate(COLS):
    x = ML + i * (cw + 0.31)
    line(s, x, BODY_Y, cw, col, 2.0)
    write(textbox(s, x, BODY_Y + 0.16, cw, 0.34).paragraphs[0], head, 14.5,
          INK, bold=True)
    prose(s, x, BODY_Y + 0.60, cw, points, size=11.8, gap=11,
          tag=f"s3 col{i}", limit=5.90)

takeaway(s, "Coding is known to interact with the front end, and attribution is "
            "known to be fragile — but the two literatures have never been put "
            "in the same experiment.", accent=MUTED, tag="s3 band")


# =========================================================================
# 4 — gap and contributions
# =========================================================================
s = slide("The gap, and what this paper contributes", "Research gap", "intro")

y = eyebrow(s, ML, BODY_Y, CW * 0.40, "Where the gap lies", MP3)
bullets(s, [
    (0, "Coding effects are characterised for **accuracy**, and never for "
        "**evidence**."),
    (0, "Attribution fragility is characterised for **vision** under "
        "**adversarial** perturbation, and not for **audio** under the routine "
        "degradation that actually occurs."),
    (0, "No prior work isolates the **container** under a verified constraint "
        "of signal invariance, and none calibrates an importance shift against "
        "a **matched empty perturbation**."),
], ML, y + 0.04, CW * 0.40, size=13.5, gap=11, accent=MP3, tag="s4 gap")

px = ML + CW * 0.455
eyebrow(s, px, BODY_Y, CW * 0.545, "Contributions", WAV)

CONTRIB = [
    ("1", "A matched-format decomposition of the codec effect", "38.6 points of "
     "loss recovered by training in the served condition — the codec does not "
     "remove the affective information."),
    ("2", "Tree structure as a readable attribution", "Which named test the "
     "codec changes, at which threshold, and how far the agreement reaches; "
     "entropy at every node separates structural from informational agreement."),
    ("3", "Null calibration for feature-importance shift", "Three floors — "
     "tie-breaking, actor resampling, and an informationally empty container "
     "re-wrap — under a paired within-draw design."),
    ("4", "A repair that is diagnosable and cheap", "One named descriptor "
     "recovers 98 % of a kernel machine's loss, at no cost on clean audio."),
    ("5", "A class-level reading of the mismatch failure", "The mismatched "
     "network abandons half the label space — arithmetically invisible to "
     "balanced accuracy."),
    ("6", "Two negative results reported rather than smoothed over", "XAI "
     "methods disagree on one fitted model; a global ranking can look stable "
     "while the classifier has stopped working."),
]
y = BODY_Y + 0.32
for num, head, sub in CONTRIB:
    write(textbox(s, px, y, 0.30, 0.26).paragraphs[0], num, 13, WAV, bold=True)
    tf = textbox(s, px + 0.28, y - 0.015, CW * 0.545 - 0.28, 0.72)
    p = tf.paragraphs[0]
    p.line_spacing = 1.14
    write(p, head + ".  ", 12.5, INK, bold=True)
    write(p, sub, 11.5, MUTED)
    n = nlines(head + ".   " + sub, CW * 0.545 - 0.28, 12.0)
    check(f"s4 contrib {num}", y + n * 12.0 * 1.21 * 1.14 / 72.0)
    y += 0.82


# =========================================================================
# 5 — dataset and stimulus design
# =========================================================================
s = slide("CREMA-D, rendered in four coding conditions",
          "Dataset and experimental design", "method")

LW = CW * 0.44
y = eyebrow(s, ML, BODY_Y, LW, "The corpus", WAV)
table(s, [
    ["", ""],
    ["Clips", "7,442  (7,441 usable)"],
    ["Actors", "91  —  48 male, 43 female, ages 20–74"],
    ["Sentences", "12 fixed, identical across emotions"],
    ["Emotions", "6  —  ANG  DIS  FEA  HAP  NEU  SAD"],
    ["Clip length", "2.54 s ± 0.51"],
    ["Split", "speaker-independent, 60 / 11 / 20 actors"],
    ["", "4,905 / 896 / 1,640 clips"],
], ML, y, LW, col_w=[0.30, 0.70], size=11.5, head_h=0.06, row_h=0.245,
   align=["l", "l"])

bullets(s, [
    (0, "Holding the **words constant** removes the lexical channel. Any signal "
        "a model recovers is carried by acoustic realisation — which is exactly "
        "what perceptual coding operates on."),
    (0, "Human voice-only ceiling is **45.5 %** (NEU recall 0.966, SAD 0.164) — "
        "context for the absolute accuracies, not a target."),
    (0, "Speaker η² = 0.178 against emotion η² = 0.098, and **340 of 436** "
        "descriptors are speaker-dominant. Speaker-independent splitting is "
        "load-bearing here, not a refinement."),
    (0, "Extraction failed on exactly one clip, `1076_MTI_SAD_XX` — the one the "
        "corpus documents as containing no audio."),
], ML, y + 1.84, LW, size=11.8, gap=9, accent=WAV, tag="s5 left")

px, RW = ML + CW * 0.485, CW * 0.515
y = eyebrow(s, px, BODY_Y, RW, "The four conditions", WAV)
y = prose(s, px, y, RW, [
    "Canonical reference: 16 kHz mono 16-bit PCM, loudness-normalised to "
    "EBU R128 (I = −23 LUFS, LRA 7, TP −2 dB). One encoder invocation per "
    "condition; one identical FFmpeg call decodes them all, so the decoder "
    "never confounds the codec.",
], size=11.5, color=MUTED, tag="s5 ref")

table(s, [
    ["Condition", "What it is", "What it varies"],
    ["ref", "uncompressed WAV", "the reference"],
    ["mp3_64", "MP3 @ 64 kbit/s", "codec"],
    ["mp4_aac64", "MP4/AAC @ 64 kbit/s", "codec"],
    ["roundtrip_wav", "AAC bitstream re-wrapped as WAV", "container only"],
], px, y + 0.14, RW, col_w=[0.27, 0.44, 0.29], size=11.5, head_h=0.30,
   row_h=0.30, align=["l", "l", "l"], rule_rows=(4,))

py = y + 1.68
panel(s, px, py, RW, 1.86, TINT)
line(s, px, py, 0.055, MP3, 1.86 * 72)
prose(s, px + 0.28, py + 0.16, RW - 0.52, [
    ("`roundtrip_wav` is the device that makes the study identifiable.",
     12.5, INK),
    "It carries the **same codec output** through a different container, so it "
    "changes the decode path and no coding information at all. Median "
    "standardised feature difference **|SMD| = 7.6 × 10⁻⁵**; only 12 of 436 "
    "descriptors exceed 0.05. That is the noise floor of this study, and it is "
    "used as an **informationally empty perturbation** rather than merely "
    "reported as small.",
], size=11.5, gap=7, tag="s5 panel", limit=py + 1.72)


# =========================================================================
# 6 — feature representation
# =========================================================================
s = slide("436 named acoustic descriptors, not learned features",
          "Methodology  ·  Representation", "method")

LW = CW * 0.585
table(s, [
    ["Family", "n", "Contents"],
    ["Cepstral", "240", "20 MFCCs with their Δ and ΔΔ temporal derivatives"],
    ["Spectral", "114", "centroid, bandwidth, roll-off 85/95 %, flatness, flux, "
                        "entropy, slope, ZCR, RMS, 7-band contrast, 8 band-energy "
                        "ratios, HF survival above 4 and 6 kHz"],
    ["Prosodic", "56", "F₀ statistics and slope, jitter (local/RAP/PPQ5/DDP), "
                       "shimmer (local/dB/APQ3/APQ5/APQ11/DDA), HNR, formants "
                       "F₁–F₅ with bandwidths, intensity, voiced timing, CPPS"],
    ["Chroma + global", "26", "12 chroma bins; clip duration and peak amplitude"],
], ML, BODY_Y, LW, col_w=[0.20, 0.08, 0.72], size=11.3, head_h=0.30,
   row_h=0.62, align=["l", "c", "l"])

bullets(s, [
    (0, "Every frame-level contour is collapsed to **eight order statistics** — "
        "mean, std, min, max, median, IQR, skewness, excess kurtosis."),
    (0, "librosa with a 512-point FFT, 400-sample window, 160-sample hop; "
        "prosody in Praat through Parselmouth. Chroma tuning is **pinned to "
        "zero** — an estimated tuning would itself shift under compression and "
        "become part of the effect being measured."),
], ML, BODY_Y + 3.02, LW, size=11.8, gap=9, accent=WAV, tag="s6 left")

px, RW = ML + CW * 0.625, CW * 0.375
panel(s, px, BODY_Y, RW, 2.52, TINT)
prose(s, px + 0.28, BODY_Y + 0.18, RW - 0.52, [
    ("WHY NAMED DESCRIPTORS", 10, MP3),
    "An attribution on a latent dimension is not a claim about anything an "
    "acoustician would recognise.",
    "An attribution on `prosody_cpps` is a claim about **breathiness**. A split "
    "on `spectral_contrast_6_mean` is a claim about the **6–8 kHz band**. "
    "Everything here is exact rather than estimated, and can be checked against "
    "an acoustic reading.",
], size=12, gap=9, tag="s6 panel", limit=BODY_Y + 2.38)

y = eyebrow(s, px, BODY_Y + 2.72, RW, "Separability", WAV)
bullets(s, [
    (0, "**397 of 436** descriptors separate emotion at Bonferroni-corrected "
        "significance."),
    (0, "Strongest: `mfcc_d1_00_std` (F = 1194), `mfcc_00_std` (1159), "
        "`prosody_intensity_std` (1078) — all **variability** measures. Emotion "
        "in acted speech lives in the dynamics."),
    (0, "51 principal components for 95 % of variance: genuinely "
        "high-dimensional, not a few latent factors under 436 names."),
], px, y, RW, size=11.5, gap=8, accent=WAV, tag="s6 right")


# =========================================================================
# 7 — models and the splitting criterion
# =========================================================================
s = slide("Fifteen models, an explicit splitting criterion, ten XAI methods",
          "Methodology  ·  Models", "method")

LW = CW * 0.375
y = eyebrow(s, ML, BODY_Y, LW, "The zoo", WAV)
y = bullets(s, [
    (0, "**14 scikit-learn classifiers** spanning tree ensembles, kernel "
        "machines, linear and discriminant models, and instance-based and "
        "probabilistic baselines."),
    (0, "**1 PyTorch MLP** — 512-256-128, batch norm, dropout 0.3, AdamW under "
        "a cosine schedule, class-weighted cross-entropy with label smoothing, "
        "early stopping on validation balanced accuracy."),
], ML, y, LW, size=12, gap=10, accent=WAV, tag="s7 zoo")

y = eyebrow(s, ML, y + 0.28, LW, "The XAI stack", WAV)
bullets(s, [
    (0, "**Tabular** — TreeSHAP (exact for ensembles), LinearSHAP, KernelSHAP, "
        "LIME, permutation importance, partial dependence, and a depth-4 global "
        "surrogate tree."),
    (0, "**Network** — six Captum methods: Saliency, Input×Gradient, Integrated "
        "Gradients, DeepLIFT, GradientSHAP, Feature Ablation."),
    (0, "All attribution is measured on the **held-out speakers**."),
], ML, y, LW, size=12, gap=10, accent=WAV, tag="s7 xai")

px, RW = ML + CW * 0.425, CW * 0.575
panel(s, px, BODY_Y - 0.04, RW, 5.28, TINT)
y = eyebrow(s, px + 0.34, BODY_Y + 0.16, RW - 0.68, "Entropy and information gain", MP3)
y = prose(s, px + 0.34, y, RW - 0.68, [
    "Trees are grown with the **Shannon-entropy criterion** rather than Gini, "
    "so the split rules reported later are exactly the rules these equations "
    "describe. Class weights are balanced.",
], size=12, tag="s7 eq intro")

tf = textbox(s, px + 0.34, y + 0.14, RW - 0.68, 0.94)
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
for t, sz in [("H(S)  =  −  Σ", 17), ("c", 11), ("   p", 17), ("c", 11),
              ("   log", 17), ("2", 11), ("   p", 17), ("c", 11)]:
    write(p, t, sz, INK)
p = tf.add_paragraph()
p.space_before = Pt(13)
p.alignment = PP_ALIGN.CENTER
for t, sz in [("IG(S, f, τ)  =  H(S)  −  (w", 17), ("L", 11), ("/w) H(S", 17),
              ("L", 11), (")  −  (w", 17), ("R", 11), ("/w) H(S", 17),
              ("R", 11), (")", 17)]:
    write(p, t, sz, INK)

line(s, px + 0.34, y + 1.20, RW - 0.68, RULE, 0.75)
y2 = eyebrow(s, px + 0.34, y + 1.36, RW - 0.68,
             "Worked example  ·  the root of every tree", MP3)
bullets(s, [
    (0, "Balanced weights make the root entropy identical in all four "
        "conditions: **H = log₂ 6 = 2.5850 bits**."),
    (0, "On uncompressed audio the maximising test is "
        "`prosody_intensity_std ≤ 8.2184`, sending 3098.67 of the 4905.00 "
        "weighted samples left."),
    (0, "IG = 2.5850 − 1.5066 − 0.7693 = **0.3090 bits**."),
    (0, "The single most informative question about a CREMA-D clip is *how much "
        "does its loudness vary*, and the answer resolves **12 %** of the total "
        "label uncertainty."),
], px + 0.34, y2, RW - 0.68, size=12, gap=8, accent=MP3,
   tag="s7 worked", limit=BODY_Y + 5.06)


# =========================================================================
# 8 — null calibration protocol
# =========================================================================
s = slide("An attribution-shift claim needs a null: the ranking itself is fragile",
          "Methodology  ·  Null calibration", "method")

LW = CW * 0.46
y = bullets(s, [
    (0, "A claim of the form *“compression changed which features the model "
        "relies on”* is easy to produce and hard to justify."),
    (0, "Feature-importance rankings are unstable by construction. Two "
        "descriptors that split the data almost equally well trade places, and "
        "**a single swap near the root reroutes every node below it**."),
    (0, "The general warning is in the literature. The specific discipline it "
        "implies — measure the shift induced by a manipulation that carries "
        "**no information** — is not yet standard practice."),
], ML, BODY_Y, LW, size=13, gap=11, accent=WAV, tag="s8 left")

y = eyebrow(s, ML, y + 0.30, LW, "Three floors, all on the same estimator", MP3)
table(s, [
    ["Floor", "What it varies"],
    ["Tie-breaking", "the estimator only, data held fixed"],
    ["Actor bootstrap", "which actors are fitted, codec fixed"],
    ["Container re-wrap", "the decode path, codec output fixed"],
], ML, y, LW, col_w=[0.36, 0.64], size=11.8, head_h=0.30, row_h=0.32,
   align=["l", "l"])

px, RW = ML + CW * 0.515, CW * 0.485
y = eyebrow(s, px, BODY_Y, RW, "The paired within-draw design", WAV)
for lab, txt in [
    ("Draw", "On each of **20 draws**, sample **51 of the 60** training actors "
             "without replacement. No actor is duplicated."),
    ("Floor", "Fit the tree **twice on the uncompressed audio of those actors** "
              "under different tie-breaking. That is the floor for this draw."),
    ("Effect", "Fit once on the **coded audio of the same actors**. That is the "
               "effect for this draw."),
    ("Test", "Wilcoxon signed-rank on the paired per-draw difference. Sampling "
             "variance is held fixed by construction."),
]:
    panel(s, px, y, RW, 0.72, TINT)
    line(s, px, y, 0.05, WAV, 0.72 * 72)
    write(textbox(s, px + 0.26, y + 0.06, 0.90, 0.30).paragraphs[0], lab, 12,
          WAV, bold=True)
    prose(s, px + 1.14, y + 0.06, RW - 1.40, [txt], size=11.5,
          tag=f"s8 {lab}", limit=y + 0.66)
    y += 0.86

takeaway(s, "Made **marginally**, the actor-resampling floor is wide enough to "
            "swallow every codec effect, and we would have reported nothing at "
            "all. Floor and effect must be measured inside the same draw.",
         accent=WAV, h=0.72, tag="s8 band")


# =========================================================================
# 9 — signal-level manipulation
# =========================================================================
s = slide("Two codecs, two entirely different mechanisms",
          "Results  ·  What coding does to the descriptors")

picture(s, "fig_signal", ML, BODY_Y - 0.06, CW, 2.36,
        caption="Measured on the feature table over 7,441 clips, before any model is involved. "
                "(a) the six most displaced descriptors under each codec  ·  (b) ECDF of |SMD| "
                "over all 436  ·  (c) the encoder-priming shift.")

by, cw2 = BODY_Y + 2.84, (CW - 0.34) / 2
for i, (head, col, pts) in enumerate([
    ("MP3 @ 64 kbit/s abandons the top octave", MP3, [
        "Perceptual bit allocation quantises the highest band coarsely and sets "
        "masked bins to zero.",
        "`spectral_contrast_6_mean` (the 6–8 kHz band) moves **16.38 → 46.54**, "
        "which is **38.6 training standard deviations**. Spectral flatness "
        "collapses with it, mean 0.0189 → 0.0066.",
        "The descriptor arrives at inference in a range no model fitted on clean "
        "audio has ever seen. **Concentrated damage.**",
    ]),
    ("MP4/AAC shifts frame alignment", WAV, [
        "Encoder priming lengthens every clip uniformly by **+32.68 ms**. MP3 "
        "adds only 0.15 ms.",
        "That perturbs every frame-aligned temporal derivative a little: the "
        "largest AAC shifts are all delta-MFCC means, led by `mfcc_d1_01_mean` "
        "at SMD = −4.16.",
        "**Diffuse damage.** Caveat: delta means sit near zero with small "
        "standard deviations, so their SMDs are inflated; the duration "
        "measurement is the harder evidence.",
    ])]):
    x = ML + i * (cw2 + 0.34)
    panel(s, x, by, cw2, 2.36, TINT)
    line(s, x, by, 0.05, col, 2.36 * 72)
    write(textbox(s, x + 0.26, by + 0.14, cw2 - 0.50, 0.30).paragraphs[0],
          head, 13, INK, bold=True)
    prose(s, x + 0.26, by + 0.48, cw2 - 0.50, pts, size=10.8, gap=6,
          tag=f"s9 panel{i}", limit=by + 2.24)


# =========================================================================
# 10 — in-format leaderboard
# =========================================================================
s = slide("On clean audio, intrinsic interpretability costs a third of the performance",
          "Results  ·  In-format performance")

LW = CW * 0.475
table(s, [
    ["Model", "Acc.", "Bal. acc.", "Macro F₁", "κ"],
    ["ANN (PyTorch)", "0.599", "0.599", "0.596", "0.518"],
    ["LightGBM", "0.590", "0.591", "0.584", "0.508"],
    ["SVM (RBF)", "0.585", "0.586", "0.581", "0.502"],
    ["LDA (shrinkage)", "0.585", "0.585", "0.580", "0.502"],
    ["Linear SVM", "0.583", "0.584", "0.577", "0.500"],
    ["XGBoost", "0.581", "0.582", "0.575", "0.497"],
    ["MLP (scikit-learn)", "0.573", "0.574", "0.569", "0.488"],
    ["Logistic regression", "0.565", "0.567", "0.561", "0.478"],
    ["Random forest", "0.532", "0.535", "0.517", "0.439"],
    ["kNN (k = 15)", "0.477", "0.479", "0.466", "0.373"],
    ["Decision tree", "0.390", "0.391", "0.382", "0.268"],
    ["Dummy (majority)", "0.171", "0.167", "0.049", "0.000"],
], ML, BODY_Y, LW, col_w=[0.40, 0.15, 0.17, 0.16, 0.12], size=11.5,
   head_h=0.30, row_h=0.285, rule_rows=(11, 12))

prose(s, ML, BODY_Y + 3.72, LW, [
    "Six-way task, 20 held-out speakers. Chance = 0.167; human voice-only "
    "ceiling = 0.455. Eleven of fifteen models shown.",
], size=10, color=MUTED, tag="s10 note")

px, RW = ML + CW * 0.525, CW * 0.475
for i, (v, lab, col) in enumerate([
        ("0.391", "decision tree\nbalanced accuracy", MP3),
        ("0.599", "PyTorch ANN\nbalanced accuracy", WAV),
        ("−34 %", "the price of a model\na human can read", INK)]):
    stat(s, px + i * RW / 3, BODY_Y - 0.06, RW / 3, v, lab, color=col, vsize=27)

S10_BAND = ("Every model beats the **45.5 %** human voice-only ceiling: they "
            "are trained on the *intended* label and exploit acting cues "
            "listeners do not consciously decode.")
bullets(s, [
    (0, "The decision tree is the **one model in the zoo a human can read from "
        "beginning to end**, and it reaches 0.391 against 0.599 for the network. "
        "Intrinsic interpretability costs a third of what is achievable here."),
    (0, "That argues **against Rudin's prescription** in this specific setting, "
        "and in favour of applying post-hoc attribution to a strong model."),
    (0, "The spread between the top eight models is **less than three points**, "
        "so the leaderboard ordering carries little information. The differences "
        "that matter here are **10 to 38 points**, and they appear only when the "
        "coding condition changes."),
], px, BODY_Y + 1.24, RW, size=13, gap=12, accent=MP3,
   tag="s10 right", limit=BODY_B - band_h(S10_BAND, RW, 12.5) - 0.16)

takeaway(s, S10_BAND, x=px, w=RW, size=12.5, tag="s10 band")


# =========================================================================
# 11 — train/serve mismatch
# =========================================================================
s = slide("Robustness is decided by how a model consumes a value",
          "Results  ·  Train / serve mismatch")

LW = CW * 0.52
y = eyebrow(s, ML, BODY_Y, LW,
            "axis-aligned partitioning  ·  magnitude-sensitive", MUTED, 9.5, 0.8)
table(s, [
    ["Model", "ref", "mp3_64", "Δ MP3", "aac64", "Δ AAC", "r-trip"],
    ["LightGBM", "0.591", "0.586", "−0.005", "0.559", "−0.032", "0.562"],
    ["XGBoost", "0.582", "0.586", "+0.004", "0.560", "−0.022", "0.560"],
    ["HistGradientBoosting", "0.573", "0.559", "−0.014", "0.530", "−0.043", "0.533"],
    ["Random forest", "0.535", "0.532", "−0.003", "0.522", "−0.013", "0.522"],
    ["Extra trees", "0.519", "0.521", "+0.002", "0.512", "−0.007", "0.507"],
    ["Decision tree", "0.391", "0.382", "−0.009", "0.367", "−0.023", "0.371"],
    ["kNN", "0.479", "0.362", "−0.118", "0.467", "−0.012", "0.457"],
    ["Gaussian NB", "0.430", "0.277", "−0.152", "0.437", "+0.008", "0.433"],
    ["Linear SVM", "0.584", "0.405", "−0.179", "0.471", "−0.113", "0.462"],
    ["MLP (scikit-learn)", "0.574", "0.383", "−0.191", "0.491", "−0.083", "0.494"],
    ["ANN (PyTorch)", "0.599", "0.354", "**−0.246**", "0.513", "−0.087", "0.512"],
    ["LDA (shrinkage)", "0.585", "0.338", "−0.247", "0.507", "−0.079", "0.502"],
    ["Logistic regression", "0.567", "0.233", "−0.334", "0.472", "−0.095", "0.470"],
    ["SVM (RBF)", "0.586", "0.203", "**−0.383**", "0.499", "−0.087", "0.497"],
], ML, y, LW, col_w=[0.28, 0.12, 0.13, 0.13, 0.115, 0.115, 0.11],
   size=10.8, head_h=0.28, row_h=0.272, rule_rows=(7,))

prose(s, ML, y + 4.10, LW, [
    "All fifteen models, trained on uncompressed audio only. Balanced accuracy "
    "on 20 held-out speakers.",
], size=10, color=MUTED, tag="s11 note")

px, RW = ML + CW * 0.565, CW * 0.435
S11_BAND = ("**MP4/AAC costs everyone something and nobody very much** — 0.7 to "
            "11 points regardless of family.")
bullets(s, [
    (0, "The split is **categorical, not graded**. Every tree-based model stays "
        "within ±0.014 of its clean score on MP3. Every model that consumes a "
        "descriptor as a continuous magnitude loses **12 to 38 points**."),
    (0, "The mechanism is readable because the descriptors are named. A tree "
        "asks whether `spectral_contrast_6_mean` exceeds a threshold near 17, "
        "and the answer is the same at 16.4 or at 46.5. A kernel model "
        "multiplies that **38σ excursion** straight into its decision function."),
    (0, "The groups are named for the property that separates them, not for a "
        "model family: kNN computes distances and Gaussian NB per-descriptor "
        "likelihoods. What all eight share is that **magnitude enters "
        "continuously**."),
    (0, "**Container routing has no material effect.** Across the fifteen "
        "models the median difference between `roundtrip` and `mp4_aac64` is "
        "0.003 and the maximum 0.010 — one to two orders of magnitude below the "
        "codec effects."),
], px, BODY_Y, RW, size=12.2, gap=11, accent=MP3, tag="s11 right", limit=BODY_B - band_h(S11_BAND, RW, 12.5) - 0.16)

takeaway(s, S11_BAND, x=px, w=RW, size=12.5, tag="s11 band")


# =========================================================================
# 12 — the per-class collapse
# =========================================================================
s = slide("Balanced accuracy is arithmetically unable to show what actually failed",
          "Results  ·  The collapse, opened by class")

LW = CW * 0.55
y = eyebrow(s, ML, BODY_Y, LW, "Clips the network predicts into each class, of 1,640", MP3)
table(s, [
    ["Condition", "ANG", "DIS", "FEA", "HAP", "NEU", "SAD"],
    ["ref", "299", "221", "322", "253", "248", "297"],
    ["mp3_64", "375", "**0**", "**934**", "331", "**0**", "**0**"],
    ["mp4_aac64", "319", "432", "315", "292", "84", "198"],
    ["roundtrip_wav", "293", "441", "311", "308", "82", "205"],
], ML, y, LW, col_w=[0.28, 0.12, 0.12, 0.12, 0.12, 0.12, 0.12], size=12,
   head_h=0.30, row_h=0.32)

y = eyebrow(s, ML, y + 1.68, LW, "Per-class F₁", MP3)
table(s, [
    ["Condition", "ANG", "DIS", "FEA", "HAP", "NEU", "SAD"],
    ["ref", "0.770", "0.467", "0.555", "0.567", "0.611", "0.607"],
    ["mp3_64", "0.660", "**0.000**", "0.407", "0.429", "**0.000**", "**0.000**"],
    ["mp4_aac64", "0.718", "0.416", "0.555", "0.538", "0.340", "0.481"],
    ["roundtrip_wav", "0.705", "0.422", "0.562", "0.531", "0.342", "0.495"],
], ML, y, LW, col_w=[0.28, 0.12, 0.12, 0.12, 0.12, 0.12, 0.12], size=12,
   head_h=0.30, row_h=0.32)

px, RW = ML + CW * 0.59, CW * 0.41
bullets(s, [
    (0, "Under MP3 the network **never predicts disgust, neutrality or "
        "sadness**. Not rarely — never: zero predictions across all 1,640 test "
        "clips. What survives is a three-way classifier that routes **934 "
        "clips, 57 % of the test set, into fear**."),
    (0, "The arithmetic that hides this is worth stating. Three of six "
        "per-class recalls are zero by construction, so the highest score this "
        "degenerate model could reach is **0.500** even with perfect recall on "
        "the classes it still emits. It reaches 0.354 of that — which reads as "
        "moderate degradation rather than structural failure."),
    (0, "**Macro F₁ is the more honest instrument** and falls 0.596 → 0.249; "
        "top-2 accuracy falls 0.805 → 0.477."),
    (0, "AAC shows the same mechanism mildly, which is what makes it a "
        "mechanism and not an artefact of one extreme setting: neutrality "
        "nearly vanishes (248 → 84 predictions, F₁ 0.611 → 0.340) and disgust "
        "absorbs the displaced mass (221 → 432). The container control "
        "reproduces it to within a few clips per class."),
], px, BODY_Y, RW, size=12, gap=10, accent=MP3, tag="s12 right", limit=5.98)

takeaway(s, "A monitor watching balanced accuracy alone records a serious but "
            "ordinary degradation. A monitor watching the **predicted-class "
            "histogram** — which costs nothing to compute — sees half the label "
            "space disappear.", h=0.72, tag="s12 band")


# =========================================================================
# 13 — matched-format training
# =========================================================================
s = slide("The codec does not remove the affective information",
          "Results  ·  Matched-format training")

LW = CW * 0.475
y = eyebrow(s, ML, BODY_Y, LW,
            "Balanced accuracy when fitted and evaluated in the same condition",
            MUTED, 9.5, 0.8)
table(s, [
    ["Model", "WAV", "MP3", "AAC", "r-trip", "range"],
    ["LightGBM", "0.591", "0.583", "0.582", "0.585", "0.008"],
    ["SVM (RBF)", "0.586", "**0.589**", "0.572", "0.572", "0.017"],
    ["LDA (shrinkage)", "0.585", "0.586", "0.579", "0.584", "0.007"],
    ["MLP (scikit-learn)", "0.574", "0.566", "0.569", "0.558", "0.017"],
    ["Logistic regression", "0.567", "0.563", "0.570", "0.567", "0.008"],
    ["Random forest", "0.526", "0.519", "0.524", "0.523", "0.007"],
    ["Decision tree (d = 12)", "0.408", "0.379", "0.398", "0.369", "0.039"],
    ["Decision tree (d = 3)", "0.407", "0.408", "0.389", "0.389", "0.019"],
], ML, y, LW, col_w=[0.32, 0.14, 0.14, 0.14, 0.14, 0.12], size=11.5,
   head_h=0.30, row_h=0.30)

for i, (v, lab) in enumerate([("+38.6", "points recovered\nby the RBF-SVM"),
                              ("+33.0", "logistic\nregression"),
                              ("+24.8", "LDA")]):
    stat(s, ML + i * LW / 3, y + 2.86, LW / 3, v, lab, color=VOICE, vsize=27)

px, RW = ML + CW * 0.525, CW * 0.475
bullets(s, [
    (0, "**The codec does not remove the information.** The RBF-SVM scores "
        "0.589 on MP3 when it is *trained* on MP3, against 0.586 on clean "
        "audio — a recovery of 38.6 points over its mismatched 0.203. If "
        "64 kbit/s MP3 destroyed the acoustic correlates of emotion, no model "
        "fitted on MP3 could recover them."),
    (0, "**No model has a spread greater than 0.039** across the four "
        "conditions. The catastrophic cross-format column measures an artefact "
        "of deployment practice, not a property of the codec."),
    (0, "**The tree family's advantage exists only under mismatch.** Under "
        "matched training LightGBM leads the RBF-SVM by 0.008 on clean audio "
        "and trails it by 0.006 on MP3."),
    (0, "**The fix is format-specific, not a general immunisation.** An "
        "RBF-SVM trained on `mp4_aac64` and served MP3 falls to 0.167 — exactly "
        "chance. Matched training protects against the codec it saw and no "
        "other, which is why the descriptor-level repair is the more practical "
        "remedy when the delivery format is unknown in advance."),
], px, BODY_Y, RW, size=12.2, gap=11, accent=VOICE, tag="s13 right")


# =========================================================================
# 14 — per-condition tree structure
# =========================================================================
s = slide("The same root question in every condition, then a substitution",
          "Results  ·  Tree structure as a readable attribution")

picture(s, "fig_trees", ML, BODY_Y + 0.16, CW * 0.50, 4.60)

px, RW = ML + CW * 0.525, CW * 0.475
y = eyebrow(s, px, BODY_Y, RW,
            "root  (H = 2.5850 bits)  ·  then depth 1, branch L", MUTED, 9.5, 0.8)
table(s, [
    ["Condition", "Descriptor", "Threshold", "IG"],
    ["ref", "prosody_intensity_std", "8.2184", "0.3090"],
    ["mp3_64", "prosody_intensity_std", "8.2210", "0.3090"],
    ["mp4_aac64", "prosody_intensity_std", "8.3241", "0.3064"],
    ["roundtrip_wav", "prosody_intensity_std", "8.3241", "0.3064"],
    ["ref", "mfcc_d1_00_std", "3.9965", "0.1415"],
    ["mp3_64", "mfcc_d1_00_std", "4.0175", "0.1398"],
    ["mp4_aac64", "**mfcc_00_iqr**", "60.5853", "0.1408"],
    ["roundtrip_wav", "**mfcc_00_iqr**", "60.5802", "0.1408"],
], px, y, RW, col_w=[0.28, 0.38, 0.18, 0.16], size=10.8, head_h=0.28,
   row_h=0.275, align=["l", "l", "r", "r"], rule_rows=(5,))

bullets(s, [
    (0, "**The primary evidence does not change**, and the threshold moves in "
        "step with the mechanism. MP3 shifts it by 0.03 %, AAC by 1.29 %: AAC "
        "priming adds 32.68 ms of near-silent lead-in to every clip, raising "
        "the dispersion of the intensity contour, and the tree compensates."),
    (0, "**The trees diverge at depth 1, under AAC and not under MP3.** WAV and "
        "MP3 ask about `mfcc_d1_00_std`, a *temporal derivative*; AAC abandons "
        "it for `mfcc_00_iqr`, a *static* dispersion measure of the same "
        "coefficient, at 0.1408 bits against 0.1415 — exactly what a "
        "frame-alignment perturbation should do."),
    (0, "**The container control reproduces structurally.** `mp4_aac64` and "
        "`roundtrip_wav` select identical descriptors at all seven nodes."),
], px, y + 2.58, RW, size=10.6, gap=8, accent=MP3, tag="s14 right")


# =========================================================================
# 15 — depth profile and the information budget
# =========================================================================
s = slide("Structure is lost with depth. The information budget is not.",
          "Results  ·  The entropy calculation through the whole tree")

LW = CW * 0.42
y = eyebrow(s, ML, BODY_Y, LW,
            "Descriptor agreement at path-aligned nodes", MUTED, 9.5, 0.8)
table(s, [
    ["Comparison", "d0", "d1", "d2", "d3", "d4", "d5"],
    ["Tie-breaking", "1/1", "2/2", "4/4", "8/8", "16/16", "29/30"],
    ["Empty container re-wrap", "1/1", "2/2", "4/4", "8/8", "11/16", "12/23"],
    ["MP3 @ 64k", "1/1", "2/2", "4/4", "4/8", "5/16", "4/27"],
    ["MP4/AAC @ 64k", "1/1", "1/2", "2/4", "0/8", "1/16", "0/24"],
    ["Roundtrip WAV", "1/1", "1/2", "2/4", "0/8", "0/16", "0/23"],
], ML, y, LW, col_w=[0.36, 0.10, 0.10, 0.10, 0.10, 0.12, 0.12], size=10.5,
   head_h=0.28, row_h=0.29, rule_rows=(3,))

prose(s, ML, y + 1.78, LW, [
    "Rows above the rule carry no coding information. Cells give the number of "
    "path-aligned nodes selecting the same descriptor, of those internal in "
    "both trees.",
], size=9.5, color=MUTED, tag="s15 note")

bullets(s, [
    (0, "Trees grown to **depth 6** per condition and aligned on the **path** "
        "from the root — the only identifier that means the same thing in two "
        "trees fitted on different data."),
    (0, "**Refitting on identical audio** leaves every node to depth 4 "
        "untouched, and an **informationally empty re-wrap costs three "
        "levels**: exact to depth 3, then 69 % at depth 4."),
    (0, "**Both codecs fall below that floor**, in the order the mechanisms "
        "predict: MP3 to depth 2, AAC to depth 0. Over 20 paired draws every "
        "gap is significant (p ≤ 3.5 × 10⁻⁴)."),
], ML, y + 2.02, LW, size=10.6, gap=7, accent=MP3, tag="s15 left", limit=5.96)

px, RW = ML + CW * 0.45, CW * 0.26
y2 = eyebrow(s, px, BODY_Y, RW, "Information gain by depth", MUTED, 9.5, 0.8)
table(s, [
    ["Depth", "ref", "mp3", "aac", "r-trip"],
    ["0", "0.3090", "0.3090", "0.3064", "0.3064"],
    ["1", "0.1556", "0.1546", "0.1462", "0.1462"],
    ["2", "0.0938", "0.0935", "0.0939", "0.0941"],
    ["3", "0.0832", "0.0842", "0.0863", "0.0860"],
    ["4", "0.1047", "0.1027", "0.1036", "0.1031"],
    ["5", "0.1292", "0.1275", "0.1226", "0.1220"],
    ["total", "**0.8754**", "**0.8714**", "**0.8590**", "**0.8578**"],
], px, y2, RW, col_w=[0.20, 0.20, 0.20, 0.20, 0.20], size=10.2, head_h=0.28,
   row_h=0.272, align=["l", "r", "r", "r", "r"], rule_rows=(7,))

y2 = eyebrow(s, px, y2 + 2.24, RW, "Mean node entropy H̄", MUTED, 9.5, 0.8)
table(s, [
    ["Depth", "ref", "mp3", "aac", "r-trip"],
    ["0", "2.5850", "2.5850", "2.5850", "2.5850"],
    ["2", "2.0516", "2.0506", "2.0708", "2.0708"],
    ["5", "1.7245", "1.7882", "1.7458", "1.7369"],
], px, y2, RW, col_w=[0.20, 0.20, 0.20, 0.20, 0.20], size=10.2, head_h=0.28,
   row_h=0.272, align=["l", "r", "r", "r", "r"])

px2, RW2 = ML + CW * 0.735, CW * 0.265
bullets(s, [
    (0, "Sample-weighted, so the depths are commensurable and sum to the "
        "total. Entropy falls along **almost the same trajectory in every "
        "condition**: no per-depth mean differs from the uncompressed value by "
        "more than 0.064 bits."),
    (0, "Over six levels WAV extracts **0.8754** of the 2.5850 bits available, "
        "MP3 **0.8714** (−0.46 %) and AAC **0.8590** (−1.87 %)."),
    (0, "The empty re-wrap extracts the *lowest* total of the four while "
        "carrying no coding information at all — a difference of this size sits "
        "inside the floor."),
    (0, "One caution. Gain per depth bottoms out at depth 3 and then **rises "
        "again**, 0.0832 → 0.1292 — the signature of **fitting noise**, not of "
        "resolved structure. Much of the deep tree should not be read as "
        "evidence about emotion at all."),
], px2, BODY_Y, RW2, size=11, gap=9, accent=WAV, tag="s15 right", limit=5.96)

takeaway(s, "Below depth 2 the AAC tree shares essentially no structure with the "
            "reference, and still resolves the same label uncertainty to within "
            "two parts in a hundred. **Coding changes which question the model "
            "asks, and leaves the answer worth the same.**",
         h=0.72, tag="s15 band")


# =========================================================================
# 16 — null calibration result
# =========================================================================
s = slide("Most of the apparent codec effect is generic sensitivity to perturbation",
          "Results  ·  Null calibration")

LW = CW * 0.46
y = eyebrow(s, ML, BODY_Y, LW,
            "Paired within-draw: 20 draws of 51 of 60 training actors", MUTED, 9.5, 0.8)
table(s, [
    ["", "top-25 kept", "ρ", "p vs. floor"],
    ["floor  (same audio, refit)", "**24.15**", "+0.958", "n/a"],
    ["Container re-wrap (empty)", "15.80", "+0.196", "8.1 × 10⁻⁵"],
    ["MP3 @ 64k", "12.55", "−0.068", "8.4 × 10⁻⁵"],
    ["MP4/AAC @ 64k", "10.65", "−0.137", "8.3 × 10⁻⁵"],
    ["Roundtrip WAV", "9.95", "−0.126", "7.9 × 10⁻⁵"],
], ML, y, LW, col_w=[0.42, 0.20, 0.17, 0.21], size=11.5, head_h=0.30,
   row_h=0.31, rule_rows=(2,))

picture(s, "fig_nullcal", ML, y + 2.24, LW, 2.10,
        caption="(a) made marginally, the actor-resampling floor swallows every "
                "codec effect  ·  (b) floor and effect inside the same draw.")

px, RW = ML + CW * 0.50, CW * 0.50
bullets(s, [
    (0, "**The estimator is reproducible.** With the data fixed and only "
        "tie-breaking varied, the tree keeps 24.15 of its 25 top descriptors at "
        "ρ = 0.958. Whatever else here is fragile, the fitting procedure is not."),
    (0, "**An informationally empty perturbation is expensive.** The container "
        "re-wrap is the same AAC bitstream, decoded the same way, differing by "
        "a single int16 requantisation whose median effect is 7.6 × 10⁻⁵ "
        "standardised units. It still costs **8.35 of the 25** top descriptors. "
        "A ranking that turns over by a third under a manipulation carrying no "
        "information cannot be read as a description of what the model needs."),
    (0, "**The codec effect is real, and about two fifths of its apparent "
        "size.** Measured against the empty floor, MP3 keeps 12.55 and AAC "
        "10.65. Of the 13.5 descriptors AAC appears to displace, 8.35 are "
        "generic sensitivity and only **5.15 are the codec**. Reporting the raw "
        "figure over-attributes by a factor of **2.6**."),
    (0, "**AAC displaces more of the ranking than MP3** — the reverse of the "
        "accuracy ordering, where MP3 is catastrophic and AAC mild. Coding that "
        "perturbs every temporal derivative slightly reorganises a tree more "
        "than coding that destroys one band the tree can threshold around."),
], px, BODY_Y, RW, size=12, gap=10, accent=MP3, tag="s16 right")


# =========================================================================
# 17 — XAI convergence and per-emotion evidence
# =========================================================================
s = slide("Five model families and four attribution algorithms return one answer",
          "Results  ·  Explainability, what holds")

LW = CW * 0.475
bullets(s, [
    (0, "**The convergence.** The top SHAP descriptors for XGBoost and LightGBM "
        "are identical in their first five: `mfcc_d1_00_std`, `mfcc_d2_00_std`, "
        "`duration_s`, `prosody_intensity_std`, `mfcc_00_std`. The decision tree "
        "leads with `prosody_intensity_std`, matching its own root split. All "
        "six gradient methods on the network return the same leaders."),
    (0, "The two leading descriptors are the standard deviations of the **first "
        "and second derivative of MFCC-0** — how variable the rate of change of "
        "overall loudness is, which is a direct formalisation of **vocal "
        "agitation**."),
    (0, "**Attribution mass is spread across families, not concentrated.** "
        "Integrated gradients gives MFCC 54.2 %, spectral 23.8 %, prosody "
        "15.3 %, chroma 5.7 %. MFCCs hold 55 % of the columns and take 54 % of "
        "the attribution: no per-descriptor advantage. The decision tree is the "
        "outlier, leaning on prosody at 0.306."),
], ML, BODY_Y, LW, size=12.2, gap=11, accent=VOICE, tag="s17 left", limit=5.02)

panel(s, ML, 5.16, LW, 1.58, TINT)
line(s, ML, 5.16, 0.05, VOICE, 1.58 * 72)
prose(s, ML + 0.28, 5.32, LW - 0.52, [
    ("THE ROBUST SET", 10, VOICE),
    "Only the descriptors that survive **all four independent views** should be "
    "reported as robust. Here that set is `mfcc_d1_00_std`, "
    "`prosody_intensity_std`, `duration_s`, `prosody_cpps` and "
    "`prosody_f0_median`.",
], size=12, gap=7, tag="s17 panel", limit=6.60)

px, RW = ML + CW * 0.525, CW * 0.475
y = eyebrow(s, px, BODY_Y, RW,
            "Most-attributed descriptors per emotion  (integrated gradients)",
            MUTED, 9.5, 0.8)
table(s, [
    ["", "Descriptors", "Reading"],
    ["ANG", "mfcc_d1_00_std (0.294), mfcc_d2_00_std (0.232)", "loudness dynamics"],
    ["DIS", "duration_s (0.285), prosody_f2_std (0.241)", "timing, F₂ movement"],
    ["FEA", "prosody_f0_median (0.235), prosody_hnr_std (0.181)", "raised, unstable pitch"],
    ["HAP", "duration_s (0.307), spectral_flux_std (0.247)", "spectral churn"],
    ["NEU", "mfcc_00_skew (0.327), prosody_cpps (0.327)", "steady phonation"],
    ["SAD", "prosody_cpps (0.154), mfcc_d2_00_std (0.149)", "breathiness"],
], px, y, RW, col_w=[0.10, 0.58, 0.32], size=10.5, head_h=0.28, row_h=0.40,
   align=["l", "l", "l"])

bullets(s, [
    (0, "**A construct-validity check, not a finding about emotion.** CPPS, "
        "the standard correlate of breathy phonation, leads for sadness; median "
        "F₀ leads for fear; F₂ movement leads for disgust — with no phonetic "
        "supervision."),
    (0, "**Two of these connect to the coding results.** `prosody_cpps` leads "
        "for sadness and is one of the descriptors that *leaves* the top 25 "
        "under MP3; `duration_s` leads for disgust and happiness and is what "
        "AAC priming moves by 32.68 ms. **The emotions whose evidence the "
        "codecs disturb are identifiable in advance from this table.**"),
], px, y + 2.82, RW, size=11.5, gap=9, accent=VOICE, tag="s17 right")


# =========================================================================
# 18 — XAI failure modes
# =========================================================================
s = slide("Three negative results we report rather than smooth over",
          "Results  ·  Explainability, what fails")

LW = CW * 0.42
picture(s, "fig_methods", ML, BODY_Y - 0.02, LW, 2.14,
        caption="(a) post-hoc methods on the same fitted tabular model  ·  "
                "(b) the six gradient-family methods on the network.")

bullets(s, [
    (0, "**1 · The methods disagree with each other.** SHAP tracks a model's "
        "*intrinsic* importance closely for the simpler models — ρ = 0.98 for "
        "both the decision tree and logistic regression — while LIME and "
        "permutation importance are close to unrelated to it. On the decision "
        "tree, LIME and SHAP rank descriptors at **ρ = 0.008**, while still "
        "sharing 10 of their top 20."),
    (0, "Panel (b) sharpens this rather than softening it: the six "
        "gradient-family methods agree at **ρ ≥ 0.94** for every pair. High "
        "agreement inside one family reflects a **shared mathematical "
        "lineage**, not correctness."),
    (0, "A depth-4 global surrogate reproduces its target only **57–74 %** of "
        "the time."),
], ML, BODY_Y + 2.52, LW, size=10.9, gap=9, accent=MP3, tag="s18 left")

px, RW = ML + CW * 0.455, CW * 0.545
panel(s, px, BODY_Y - 0.02, RW, 2.30, TINT)
line(s, px, BODY_Y - 0.02, 0.05, MP3, 2.30 * 72)
prose(s, px + 0.28, BODY_Y + 0.14, RW - 0.52, [
    ("2 · A BLIND SPOT IN GLOBAL ATTRIBUTION METRICS", 10.5, MP3),
    ("Under MP3, logistic regression keeps a **near-identical SHAP ranking** — "
     "ρ = 0.987, 22 of 25 top descriptors retained — while its balanced "
     "accuracy collapses from **0.567 to 0.233**.", 12.5, INK),
    "Its coefficients did not change. Only the inputs moved out of range, and a "
    "global ranking over mean |SHAP| is largely blind to that. **A global "
    "attribution-stability metric can report that nothing has changed while the "
    "classifier has effectively stopped working.**",
], size=11.4, gap=6, tag="s18 panel", limit=BODY_Y + 2.16)

y = eyebrow(s, px, BODY_Y + 2.44, RW, "3 · An estimator that failed loudly",
            MP3, 10.5, 1.0)
bullets(s, [
    (0, "KernelSHAP was applied to the RBF-SVM with **100 sampled coalitions "
        "over 436 descriptors**. It fits a weighted least-squares regression "
        "with one unknown per descriptor, so that system is rank-deficient — "
        "and it diverged. Maximum mean |SHAP| reached **2.1 × 10¹²** against a "
        "median of 5.3 × 10⁻⁴."),
    (0, "The cheap tell: a well-behaved additive attribution over 7,441 clips "
        "does not assign **whole descriptor families exactly 0.000**."),
    (0, "The ranking is excluded from every conclusion as an estimation "
        "artefact, and the implementation now **raises** when the coalition "
        "budget does not exceed the descriptor count. The SVM's trustworthy "
        "views are permutation importance and LIME, whose top descriptors "
        "reproduce the consensus."),
], px, y, RW, size=11.5, gap=8, accent=MP3, tag="s18 right")


# =========================================================================
# 19 — diagnosis and repair
# =========================================================================
s = slide("One named descriptor recovers almost the entire collapse",
          "Results  ·  Diagnosis and repair")

LW = CW * 0.44
y = eyebrow(s, ML, BODY_Y, LW,
            "Balanced accuracy on MP3 as codec-shifted descriptors are masked",
            MUTED, 9.5, 0.8)
table(s, [
    ["Masked", "SVM", "Log. reg.", "ANN", "MLP", "XGBoost", "Tree"],
    ["clean", "0.586", "0.567", "0.599", "0.574", "0.582", "0.391"],
    ["0", "0.203", "0.233", "0.354", "0.383", "0.586", "0.382"],
    ["**1**", "**0.578**", "**0.502**", "**0.580**", "**0.552**", "0.581", "0.382"],
    ["2", "0.583", "0.529", "0.587", "0.554", "0.581", "0.391"],
    ["3", "0.579", "**0.562**", "**0.595**", "0.565", "0.580", "0.391"],
    ["10", "0.575", "0.539", "0.589", "0.562", "0.582", "0.394"],
    ["80", "0.551", "0.435", "0.574", "0.523", "0.564", "0.396"],
], ML, y, LW, col_w=[0.16, 0.14, 0.18, 0.14, 0.13, 0.16, 0.13], size=11,
   head_h=0.30, row_h=0.30, rule_rows=(2,))

picture(s, "fig_neutralise", ML, y + 2.56, LW, 2.36)

px, RW = ML + CW * 0.48, CW * 0.52
S19_BAND = ("This repair has no counterpart in a system whose front end is "
            "learned: identifying the contaminated channel needed attribution "
            "over units that carry acoustic meaning.")
bullets(s, [
    (0, "Descriptors were ranked by how far each codec moves them in units of "
        "**training** standard deviation, then replaced progressively with "
        "training medians — in the compressed test set **and** the clean one."),
    (0, "The single masked descriptor is `spectral_contrast_6_mean`. Masking it "
        "alone recovers **98 % of the RBF-SVM's 38.3-point loss** and **92 % of "
        "the network's 24.6-point loss**. Three descriptors bring logistic "
        "regression to 0.562 against its 0.567 clean score."),
    (0, "**Control 1.** On uncompressed audio the identical masking costs "
        "nothing out to 20 descriptors — the network's clean control sits at "
        "0.607–0.609, marginally *above* its 0.599 baseline. These columns "
        "never carried load-bearing evidence. They were the channel through "
        "which the **encoder's fingerprint** entered the model."),
    (0, "**Control 2.** AAC behaves as its different mechanism predicts: no "
        "single descriptor dominates. The network's best recovery arrives at 10 "
        "masked (0.513 → 0.560), the SVM's at 40 (0.499 → 0.542)."),
    (0, "The diagnosis is **covariate shift on a few identifiable "
        "descriptors**, not a loss of emotional information."),
], px, BODY_Y, RW, size=11.7, gap=10, accent=VOICE, tag="s19 right",
   limit=BODY_B - band_h(S19_BAND, RW, 12.5) - 0.16)

takeaway(s, S19_BAND, accent=VOICE, x=px, w=RW, size=12.5, tag="s19 band")


# =========================================================================
# 20 — conclusion
# =========================================================================
s = slide("Three questions that a single accuracy number runs together",
          "Conclusion", "close")

LW = CW * 0.55
Q = [("1", "Does the codec destroy the information?", "No.", VOICE,
      "Matched-format training recovers 38.6 points and lands within 0.003 of "
      "the clean score."),
     ("2", "Does the codec break deployed models?", "Severely.", MP3,
      "Provided they consume descriptor magnitudes and were fitted on clean "
      "audio. And the aggregate conceals *how*: the network silently becomes a "
      "three-way classifier."),
     ("3", "Does it change what the model learns?", "Yes, but less.", AMBER,
      "The root test survives both codecs. MP3 leaves two further levels "
      "untouched, AAC none, and everything below is reorganised."),
     ("4", "Does that reorganisation cost anything?", "No.", WAV,
      "Structural and informational agreement come apart. The AAC tree shares "
      "almost no structure below depth 2 and still resolves 0.859 bits against "
      "0.875 — one question substituted for another of equal value.")]

y = BODY_Y - 0.02
for num, q, ans, col, body in Q:
    write(textbox(s, ML, y, 0.30, 0.28).paragraphs[0], num, 15, col, bold=True)
    write(textbox(s, ML + 0.30, y - 0.02, LW * 0.62, 0.30).paragraphs[0], q,
          13.5, INK, bold=True)
    write(textbox(s, ML + LW * 0.655, y - 0.02, LW * 0.35, 0.30).paragraphs[0],
          ans, 13.5, col, bold=True)
    prose(s, ML + 0.30, y + 0.28, LW - 0.30, [body], size=11.5, color=MUTED,
          tag=f"s20 q{num}", limit=y + 1.00)
    y += 1.14

px, RW = ML + CW * 0.58, CW * 0.42
panel(s, px, BODY_Y - 0.02, RW, 2.98, TINT)
line(s, px, BODY_Y - 0.02, 0.05, VOICE, 2.98 * 72)
prose(s, px + 0.28, BODY_Y + 0.14, RW - 0.52, [
    ("THE TRANSPORTABLE RESULT", 10, VOICE),
    ("An informationally empty container re-wrap already displaces **8.35 of a "
     "tree's 25** most important descriptors. Only the gap beyond that floor "
     "belongs to the codec: **3.25** for MP3 and **5.15** for AAC.", 12.5, INK),
    "We would ask reviewers to require, alongside any claim that an "
    "intervention changed a model's explanation, a measurement of the "
    "divergence induced by an intervention that carries **no information**. "
    "With it, our own claim shrinks by a factor of two and a half — and becomes "
    "defensible.",
], size=11.5, gap=8, tag="s20 panel", limit=BODY_Y + 2.84)

y = eyebrow(s, px, BODY_Y + 3.14, RW, "Limitations", MUTED)
bullets(s, [
    (0, "Acted rather than spontaneous emotion, on 12 fixed sentences. Absolute "
        "accuracies do not transfer."),
    (0, "One deterministic split, not repeated CV. The codec effects at issue "
        "are 10–38 points and are not at risk from split variance."),
    (0, "Two codecs at one bitrate; 64 kbit/s is aggressive for MP3."),
    (0, "Summary statistics discard time, so nothing here localises *when* in a "
        "clip the evidence sits."),
], px, y, RW, size=10.5, gap=7, accent=MUTED, tag="s20 lims", limit=6.72)

takeaway(s, "**Future work.** A bit-exact container null  ·  null calibration "
            "extended to SHAP and IG rankings under refitting  ·  a bitrate "
            "ladder including Opus, turning the two-point comparison into a "
            "dose–response curve.",
         accent=VOICE, x=ML, w=LW, size=12, tag="s20 band")


# =========================================================================
footer_all()
DEST.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(DEST))
print(f"wrote {DEST}  —  {len(_slides)} slides")
if WARN:
    print(f"\n{len(WARN)} layout warnings:")
    for w in WARN:
        print("  ·", w)
else:
    print("no layout warnings")
